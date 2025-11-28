"""
PowerPoint Presentation Generator Script
Run this to create the actual PowerPoint file
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()
    for line in content_lines:
        p = content.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith('  ') else 1
        p.font.size = Pt(18) if p.level == 0 else Pt(16)
    return slide

# Slide 1: Title
add_title_slide(prs,
    "AI-Powered Development: Building with Copilot & MCP",
    "A Hackathon Journey into AI-Assisted Software Engineering\nNovember 28, 2025")

# Slide 2: The Challenge
add_content_slide(prs, "The Challenge", [
    "Goal: Build full-stack task management in ONE day",
    "",
    "✅ Spring Boot 3.5.5 Backend (Java)",
    "✅ MCP Server Implementation",
    "✅ Angular UI Frontend",
    "✅ GitHub Integration",
    "✅ Full CRUD Operations",
    "",
    "The Twist: Use AI for EVERYTHING possible"
])

# Slide 3: AI Tools Arsenal
add_content_slide(prs, "AI Tools Arsenal", [
    "GitHub Copilot - Multiple Modes:",
    "  💡 Plan Mode - Brainstorming & Architecture",
    "  🤖 Agent Mode - Code Generation (GPT-4, Claude, Gemini)",
    "  💬 Chat Mode - Problem Solving & Debugging",
    "",
    "MCP Servers:",
    "  GitHub MCP - Automated repo creation, PRs, commits",
    "  Atlassian MCP - Integrated documentation",
    "  Custom MCP - Task management tools (built today!)"
])

# Slide 4: The Journey - Planning
add_content_slide(prs, "Act I: Planning with AI", [
    "Started with Copilot Plan Mode:",
    "  ✓ Brainstormed architecture patterns",
    "  ✓ Defined MCP server integration",
    "  ✓ Mapped out project structure",
    "",
    "Result: Clear roadmap in MINUTES vs. HOURS",
    "",
    "Key Learning: AI excels at structured planning"
])

# Slide 5: The Journey - Building
add_content_slide(prs, "Act II: Building with GitHub MCP", [
    "GitHub MCP Server Magic:",
    "  ✅ Created GitHub repository",
    "  ✅ Initialized project structure",
    "  ✅ Committed initial code",
    "  ✅ Created branches & PRs",
    "",
    "Copilot Agent Mode generated:",
    "  • 4 controllers • 11 MCP tools",
    "  • REST APIs • Validation & error handling",
    "",
    "Felt like: A conductor, not a coder"
])

# Slide 6: MCP Server Built
add_content_slide(prs, "Act III: MCP Server Implementation", [
    "Built Production-Grade MCP Server:",
    "",
    "📦 11 Tools - Project & task CRUD",
    "💬 3 Prompts - AI summaries & planning",
    "📚 4 Resources - Direct data access",
    "🔄 SSE Transport - Real-time streaming",
    "",
    "~2000+ lines of code (85% AI-generated)",
    "Based on dvaas pattern, adapted for tasks"
])

# Slide 7: What Worked
add_content_slide(prs, "What Worked Brilliantly ✨", [
    "🎯 The Wins:",
    "",
    "✅ Speed: 8-hour project done in 4-5 hours",
    "✅ Code Quality: Best practices built-in",
    "✅ Learning: Explored 3 AI models",
    "✅ Integration: GitHub workflow automated",
    "✅ Documentation: Comprehensive auto-docs",
    "",
    "Best Moment: AI creating PRs autonomously"
])

# Slide 8: The Challenges
add_content_slide(prs, "The Reality Check ⚠️", [
    "The Challenges:",
    "",
    "❌ Non-Deterministic: Varied results per run",
    "❌ Context Loss: Easy to lose track",
    "❌ Over-Correction: AI 'fixed' working code",
    "❌ MCP Integration: Couldn't fully integrate",
    "❌ Steering Required: Constant guidance needed",
    "",
    "AI is a powerful assistant, not autopilot (yet)"
])

# Slide 9: The Feeling
add_content_slide(prs, "The 'Just Approving' Feeling", [
    "Emotional Rollercoaster:",
    "",
    "😊 Exciting: '200 lines in seconds!'",
    "🤔 Unsettling: 'Am I even coding?'",
    "😰 Challenging: 'What did it change?'",
    "😅 Humbling: 'I'm a project manager now'",
    "",
    "Reality: Still needed deep technical knowledge",
    "  • Validate suggestions • Debug issues",
    "  • Architect solution • Keep vision coherent"
])

# Slide 10: By The Numbers
add_content_slide(prs, "By The Numbers 📊", [
    "Hackathon Stats:",
    "",
    "⏱️ Time Spent: ~6 hours",
    "📝 Lines Written: ~2,500+",
    "🤖 AI-Generated: ~85%",
    "🔨 Manual Fixes: ~15%",
    "🐛 Bugs Created by AI: 7",
    "🐛 Bugs Fixed by AI: 12",
    "",
    "ROI: 3-4x faster than traditional coding"
])

# Slide 11: Key Insights
add_content_slide(prs, "Key Insights 💡", [
    "What I Learned:",
    "",
    "1. Context is King - Better input = better output",
    "2. Iterate Rapidly - Don't aim for perfection",
    "3. Version Control is Critical - Track changes",
    "4. Trust but Verify - Always review AI output",
    "5. Multiple Models - Different strengths",
    "6. MCP is Powerful - Standard AI protocol",
    "",
    "The skill: Knowing WHAT to build > HOW to code"
])

# Slide 12: The Paradigm Shift
add_content_slide(prs, "The Future is Here 🔮", [
    "Paradigm Shift in Development:",
    "",
    "Traditional → AI-Assisted",
    "  Write every line → Orchestrate & validate",
    "  Hours on boilerplate → Minutes on structure",
    "  Manual documentation → Auto-generated",
    "  Sequential development → Parallel exploration",
    "",
    "New Skill: Knowing WHAT > HOW"
])

# Slide 13: What's Next
add_content_slide(prs, "What's Next? 🚀", [
    "Future Improvements:",
    "",
    "Technical:",
    "  • Complete MCP Copilot integration",
    "  • Add authentication & authorization",
    "  • Deploy to cloud",
    "",
    "AI Enhancements:",
    "  • Fine-tune prompts for accuracy",
    "  • Custom models on codebase",
    "  • Automated test generation"
])

# Slide 14: Closing
add_content_slide(prs, "The Takeaway", [
    "",
    "",
    "\"AI didn't replace the developer...",
    "",
    "It made me 3x more productive",
    "",
    "and 10x more exploratory\"",
    "",
    "",
    "Questions?"
])

# Save presentation
prs.save('AI_Hackathon_Presentation.pptx')
print("✅ PowerPoint created: AI_Hackathon_Presentation.pptx")

